from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "output" / "pptx"
TMP_DIR = ROOT / "tmp" / "pptx"
OUT_FILE = OUT_DIR / "STS_Investor_Business_Plan.pptx"


COLORS = {
    "bg": "0A0908",
    "navy": "0B1426",
    "surface": "141210",
    "surface2": "1C1A17",
    "line": "3D3A37",
    "text": "FAF8F5",
    "muted": "C4B9A8",
    "faint": "9A9590",
    "gold": "D4A006",
    "gold_light": "FACC15",
    "teal": "06B6D4",
    "green": "10B981",
    "amber": "F59E0B",
    "coral": "F47260",
    "white": "FFFFFF",
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_bg(slide, color: str = "bg") -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(COLORS[color])


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: str = "text",
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
    margin: float = 0.06,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_para_lines(
    slide,
    lines: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 16,
    color: str = "muted",
    bold_first: bool = False,
    gap: int = 6,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.bold = bold_first and idx == 0
        r.font.color.rgb = rgb(COLORS["text" if idx == 0 and bold_first else color])
    return box


def rect(slide, x, y, w, h, fill="surface", line="line", radius=False, transparency=0):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.fill.transparency = transparency
    shp.line.color.rgb = rgb(COLORS[line])
    shp.line.width = Pt(0.8)
    return shp


def line(slide, x1, y1, x2, y2, color="line", width=1.2):
    thickness = max(0.006, width / 120)
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(min(x1, x2)),
        Inches(y1 - thickness / 2),
        Inches(abs(x2 - x1)),
        Inches(thickness),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[color])
    shp.line.color.rgb = rgb(COLORS[color])
    shp.line.width = Pt(0)
    return shp


def circle(slide, x, y, d, fill="gold", line_color="gold"):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(COLORS[fill])
    shp.line.color.rgb = rgb(COLORS[line_color])
    return shp


def kicker(slide, label: str, idx: int) -> None:
    circle(slide, 0.55, 0.34, 0.12, "gold", "gold")
    add_text(slide, label.upper(), 0.75, 0.28, 2.7, 0.25, 8, "gold_light", True, font="Aptos", margin=0)
    add_text(slide, f"{idx:02d}", 9.05, 0.25, 0.35, 0.22, 8, "faint", True, align=PP_ALIGN.RIGHT, margin=0)
    line(slide, 9.46, 0.36, 9.75, 0.36, "line", 0.8)


def title(slide, kicker_text: str, claim: str, idx: int, subtitle: str | None = None) -> None:
    kicker(slide, kicker_text, idx)
    title_size = 24 if len(claim) > 68 else 26
    title_size = 22 if len(claim) > 88 else title_size
    add_text(slide, claim, 0.55, 0.62, 9.1, 0.78, title_size, "text", True, font="Aptos Display", margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.28, 8.75, 0.34, 10.5, "muted", False, margin=0)


def footer(slide, note="Planning assumptions for discussion. Not a forecast or guarantee.") -> None:
    line(slide, 0.55, 5.22, 9.45, 5.22, "line", 0.5)
    add_text(slide, note, 0.55, 5.3, 7.7, 0.18, 7, "faint", False, margin=0)
    add_text(slide, "Sierra Tech Spaces", 8.2, 5.3, 1.25, 0.18, 7, "gold", True, align=PP_ALIGN.RIGHT, margin=0)


def logo_path() -> Path:
    candidates = [
        ROOT / "public" / "logo-gold.png",
        ROOT / "STS_Logo_Gold_Outlined_Transparent.png",
        ROOT / "STS_Logo_Gold_Outlined.png",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("No STS logo asset found")


def add_logo(slide, x=0.58, y=0.42, w=1.02) -> None:
    path = logo_path()
    with Image.open(path) as img:
        ratio = img.height / img.width
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(w * ratio))


def card(slide, x, y, w, h, heading, body, accent="gold", value=None):
    rect(slide, x, y, w, h, "surface", "line", radius=True)
    rect(slide, x, y, 0.05, h, accent, accent)
    if value:
        add_text(slide, value, x + 0.18, y + 0.12, w - 0.35, 0.32, 20, accent, True, font="Aptos Display", margin=0)
        add_text(slide, heading, x + 0.18, y + 0.52, w - 0.35, 0.25, 12, "text", True, margin=0)
        add_text(slide, body, x + 0.18, y + 0.82, w - 0.35, h - 0.88, 9, "muted", False, margin=0)
    else:
        if h < 0.75:
            add_text(slide, heading, x + 0.18, y + 0.1, w - 0.35, 0.16, 9.5, "text", True, margin=0)
            add_text(slide, body, x + 0.18, y + 0.31, w - 0.35, max(0.12, h - 0.36), 7.5, "muted", False, margin=0)
        else:
            add_text(slide, heading, x + 0.18, y + 0.16, w - 0.35, 0.28, 13, "text", True, margin=0)
            add_text(slide, body, x + 0.18, y + 0.5, w - 0.35, max(0.12, h - 0.56), 9.5, "muted", False, margin=0)


def bar_chart(slide, x, y, w, h, rows, max_value, colors):
    label_w = 1.65
    for i, (label, value, suffix) in enumerate(rows):
        yy = y + i * (h / len(rows))
        add_text(slide, label, x, yy + 0.03, label_w, 0.22, 8.5, "muted", margin=0)
        line(slide, x + label_w, yy + 0.16, x + w, yy + 0.16, "line", 0.5)
        bw = (w - label_w - 0.42) * value / max_value
        rect(slide, x + label_w, yy + 0.06, bw, 0.18, colors[i % len(colors)], colors[i % len(colors)])
        add_text(slide, suffix, x + label_w + bw + 0.08, yy + 0.03, 0.8, 0.18, 8, "text", True, margin=0)


def create_deck() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    TMP_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_logo(s, 0.58, 0.45, 1.25)
    add_text(s, "Sierra Tech Spaces", 0.58, 1.45, 5.4, 0.36, 15, "gold_light", True, margin=0)
    add_text(s, "Investor Business Plan", 0.55, 1.83, 6.2, 1.1, 38, "text", True, font="Aptos Display", margin=0)
    add_text(s, "AI that works as hard as Egyptian businesses do.", 0.58, 2.92, 5.7, 0.4, 15, "muted", False, margin=0)
    card(s, 6.55, 0.7, 2.85, 0.82, "Launch focus", "Egyptian SMBs with manual workflows and clear ROI pains", "teal")
    card(s, 6.55, 1.72, 2.85, 0.82, "Business model", "Free audit -> focused demo -> setup fee -> monthly retainer", "gold")
    card(s, 6.55, 2.74, 2.85, 0.82, "Growth engine", "Paid ads + dedicated caller + founder-led closing", "green")
    add_text(s, "Cairo, Egypt | 2026", 0.58, 4.78, 2.3, 0.22, 9, "faint", margin=0)
    add_text(s, "Prepared for investor discussion", 6.98, 4.78, 2.4, 0.22, 9, "faint", align=PP_ALIGN.RIGHT, margin=0)

    # 2
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Problem", "Egyptian SMBs are losing margin in work they still do by hand.", 2)
    card(s, 0.65, 1.65, 2.05, 2.55, "Missed leads", "WhatsApp, Instagram, phone, and website inquiries are often handled slowly or forgotten.", "coral", "01")
    card(s, 2.95, 1.65, 2.05, 2.55, "Manual operations", "Reports, invoices, stock updates, bookings, and follow-ups depend on staff memory.", "amber", "02")
    card(s, 5.25, 1.65, 2.05, 2.55, "No visibility", "Owners cannot see response speed, task status, sales leakage, or team bottlenecks in one place.", "teal", "03")
    card(s, 7.55, 1.65, 1.8, 2.55, "Low digital maturity", "Most firms need practical systems, not abstract AI consulting.", "green", "04")
    footer(s)

    # 3
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Thesis", "AI value comes from cutting cost and speeding up work.", 3, "Adapted from the reference PDF: measurable value creation, not experimental technology.")
    card(s, 0.65, 1.55, 2.6, 1.12, "Structural tailwind", "High labor intensity, repetitive workflows, and thin margins.", "gold")
    card(s, 0.65, 2.9, 2.6, 1.12, "Best entry window", "Pre-AI workflows: spreadsheets, fragmented systems, manual follow-up.", "teal")
    bar_chart(
        s,
        3.65,
        1.55,
        5.45,
        2.5,
        [
            ("Response speed", 85, "Instant"),
            ("Manual admin", 65, "Lower"),
            ("Lead capture", 75, "Higher"),
            ("Owner visibility", 55, "Clearer"),
            ("Process cost", 45, "Reduced"),
        ],
        100,
        ["teal", "green", "gold", "amber", "coral"],
    )
    add_text(s, "STS turns this into small projects that save hours, capture revenue, and create retainers.", 3.65, 4.35, 5.55, 0.38, 11.5, "text", True, margin=0)
    footer(s, "Source logic adapted from AI-Driven Acquisition & Value Creation reference PDF; STS assumptions are local planning estimates.")

    # 4
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Targets", "We start where pain is obvious and buying cycles are short.", 4)
    targets = [
        ("E-commerce", "Missed WhatsApp leads, product uploads, abandoned carts.", "gold"),
        ("Real estate", "Lead overload, qualification, follow-up discipline.", "teal"),
        ("Medical clinics", "Bookings, repeated questions, reminders, no-shows.", "green"),
        ("Hospitality / F&B", "Reservations, reviews, menus, promotions.", "amber"),
        ("Professional services", "Client onboarding, documents, reports, follow-ups.", "coral"),
        ("Ops-heavy SMBs", "Inventory, invoices, ERP gaps, staff task tracking.", "gold"),
    ]
    for i, (h, b, c) in enumerate(targets):
        x = 0.65 + (i % 3) * 3.0
        y = 1.55 + (i // 3) * 1.45
        card(s, x, y, 2.65, 1.08, h, b, c)
    add_text(s, "Screening rule: high repetition + clear owner pain + measurable outcome + ability to demo quickly.", 0.68, 4.58, 8.5, 0.24, 11, "gold_light", True, margin=0)
    footer(s)

    # 5
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Operating model", "The client journey is designed to remove trust friction before payment.", 5)
    steps = [
        ("Free audit", "Map daily workflows and find time/money leakage.", "gold"),
        ("Workflow map", "Choose the smallest high-ROI process to modernize.", "teal"),
        ("Focused demo", "Show a working version using their real business flow.", "green"),
        ("Paid build", "Implement only what is needed, not a bloated system.", "amber"),
        ("Retainer", "Monitor, improve, train staff, and add the next automation.", "coral"),
    ]
    start_x = 0.62
    for i, (h, b, c) in enumerate(steps):
        x = start_x + i * 1.86
        circle(s, x + 0.62, 1.6, 0.46, c, c)
        add_text(s, f"{i+1}", x + 0.62, 1.71, 0.46, 0.14, 12, "bg", True, align=PP_ALIGN.CENTER, margin=0)
        if i < len(steps) - 1:
            line(s, x + 1.12, 1.83, x + 1.74, 1.83, "line", 1.2)
        card(s, x, 2.28, 1.62, 1.58, h, b, c)
    add_text(s, "Core promise: tell us what slows you down, and we turn it into a working demo before asking for a full commitment.", 0.85, 4.33, 8.3, 0.35, 12, "text", True, align=PP_ALIGN.CENTER, margin=0)
    footer(s)

    # 6
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Services", "STS sells practical systems, not vague AI transformation.", 6)
    services = [
        ("WhatsApp AI", "FAQs, lead capture, booking, routing"),
        ("Lead generation", "Ads -> qualifier -> CRM -> follow-up"),
        ("Customer service", "Multi-channel first response + handoff"),
        ("Operations automation", "Invoices, reports, stock alerts, tasks"),
        ("ERP / workflow systems", "Simple internal systems for sales, inventory, ops"),
        ("Websites", "Arabic/English lead capture sites"),
        ("Content engines", "30-day calendars, captions, brand voice"),
        ("Dashboards", "Owner visibility into leads, sales, work"),
        ("Custom software", "Portals, internal tools, client apps"),
    ]
    colors = ["gold", "teal", "green", "amber", "coral", "gold", "teal", "green", "amber"]
    for i, (h, b) in enumerate(services):
        x = 0.55 + (i % 3) * 3.05
        y = 1.42 + (i // 3) * 1.13
        card(s, x, y, 2.75, 0.86, h, b, colors[i])
    footer(s)

    # 7
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "GTM engine", "Paid demand plus calling builds the meeting machine.", 7)
    lanes = [
        ("SM agency", "FB/IG creative, lead forms, retargeting, reporting.", "gold"),
        ("Sales hire", "Calls leads, qualifies pain, books meetings.", "teal"),
        ("Founders", "Run audits, demo solutions, close paid builds.", "green"),
        ("CRM loop", "Track source, call status, proposal, retainer.", "amber"),
    ]
    for i, (h, b, c) in enumerate(lanes):
        x = 0.72 + i * 2.25
        card(s, x, 1.65, 1.85, 2.05, h, b, c)
        if i < len(lanes) - 1:
            line(s, x + 1.92, 2.62, x + 2.13, 2.62, "line", 1.5)
    add_text(s, "The founders stay focused on high-value work: audits, demos, proposals, and delivery quality.", 0.85, 4.3, 8.3, 0.3, 12, "text", True, align=PP_ALIGN.CENTER, margin=0)
    footer(s)

    # 8
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Funnel", "A small close rate can start producing revenue.", 8)
    funnel = [
        ("Ad leads", 400, "Monthly paid/social leads"),
        ("Qualified calls", 80, "20% fit after sales screening"),
        ("Meetings", 22, "Founder discovery calls"),
        ("Free audits", 14, "Workflow reviewed"),
        ("Proposals", 8, "Demo or scoped offer"),
        ("Closed projects", 3, "Setup fee + retainer path"),
    ]
    max_v = 400
    for i, (label, value, desc) in enumerate(funnel):
        y = 1.35 + i * 0.52
        bw = 5.4 * value / max_v
        rect(s, 2.35, y, bw, 0.28, ["gold", "teal", "green", "amber", "coral", "gold"][i], ["gold", "teal", "green", "amber", "coral", "gold"][i])
        add_text(s, label, 0.7, y + 0.02, 1.35, 0.18, 9, "text", True, margin=0)
        add_text(s, str(value), 8.05, y + 0.01, 0.65, 0.18, 10, "gold_light", True, align=PP_ALIGN.RIGHT, margin=0)
        add_text(s, desc, 8.82, y + 0.02, 0.82, 0.18, 7, "faint", margin=0)
    card(s, 0.72, 4.55, 2.6, 0.55, "Base monthly output", "3 closed projects can cover a lean early operating base.", "green")
    card(s, 3.65, 4.55, 2.6, 0.55, "Main constraint", "Sales follow-up discipline, not technical ability.", "amber")
    card(s, 6.58, 4.55, 2.6, 0.55, "Investor value", "Funds speed up ads, sales hiring, and repeatable demos.", "teal")
    footer(s)

    # 9
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Revenue model", "Setup fees create cash; retainers create stability.", 9)
    rows = [
        ("Quick wins", "EGP 5K-30K", "Audit fix, website, content engine, WhatsApp starter", "3-10 days"),
        ("Core systems", "EGP 20K-75K + 5K-15K/mo", "Lead gen, customer service, ops automation", "2-4 weeks"),
        ("Premium / ERP", "EGP 75K-150K+ + retainer", "Workflow system, ERP-light, dashboards, custom tools", "4-8 weeks"),
        ("Performance upside", "Selective bonus", "Bonus per qualified lead, recovered revenue, saved hours", "After proof"),
    ]
    y = 1.45
    widths = [1.55, 2.1, 3.65, 1.25]
    heads = ["Tier", "Pricing", "What it covers", "Delivery"]
    x0 = 0.6
    for j, head in enumerate(heads):
        rect(s, x0 + sum(widths[:j]), y, widths[j], 0.36, "gold", "gold")
        add_text(s, head, x0 + sum(widths[:j]) + 0.05, y + 0.08, widths[j] - 0.1, 0.12, 8, "bg", True, margin=0)
    for i, row in enumerate(rows):
        yy = y + 0.43 + i * 0.65
        for j, cell in enumerate(row):
            rect(s, x0 + sum(widths[:j]), yy, widths[j], 0.56, "surface", "line")
            add_text(s, cell, x0 + sum(widths[:j]) + 0.06, yy + 0.08, widths[j] - 0.12, 0.34, 8.2, "text" if j == 0 else "muted", j == 0, margin=0)
    add_text(s, "Pricing anchors to ROI: time saved, faster replies, recovered leads, fewer errors, and clearer owner visibility.", 0.68, 4.45, 8.45, 0.3, 12, "text", True, align=PP_ALIGN.CENTER, margin=0)
    footer(s)

    # 10
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "First 90 days", "The launch plan is fast: prove, sell, repeat.", 10)
    months = [
        ("Days 1-15", "Foundation", "Landing page, CRM, ad briefs, sales script, demo templates."),
        ("Days 16-30", "First proof", "Sprynter pilot, 2 demo builds, case study screenshots."),
        ("Days 31-60", "Outbound engine", "Ads live, caller books meetings, 20+ audits, proposals."),
        ("Days 61-90", "Scale what works", "Optimize ads, close retainers, build vertical-specific demos."),
    ]
    line(s, 1.05, 2.28, 8.9, 2.28, "line", 1.4)
    for i, (period, h, b) in enumerate(months):
        x = 0.8 + i * 2.25
        circle(s, x + 0.16, 2.09, 0.38, ["gold", "teal", "green", "amber"][i], ["gold", "teal", "green", "amber"][i])
        add_text(s, str(i + 1), x + 0.16, 2.19, 0.38, 0.12, 9, "bg", True, align=PP_ALIGN.CENTER, margin=0)
        card(s, x, 2.78, 1.78, 1.35, h, b, ["gold", "teal", "green", "amber"][i], period)
    add_text(s, "Success marker by day 90: a repeatable funnel, 2-4 paying clients, 1-2 retainers, and proof assets investors can inspect.", 0.78, 4.6, 8.55, 0.28, 11, "gold_light", True, align=PP_ALIGN.CENTER, margin=0)
    footer(s)

    # 11
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "12-month scenarios", "The base case targets EGP 1.2M revenue with lean burn.", 11, "These are planning assumptions for discussion, not guaranteed results.")
    scenarios = [
        ("Conservative", 314, 8, 3, "Proof-of-market, slower ad learning", "amber"),
        ("Base", 1188, 18, 8, "Repeatable sales motion, moderate retainers", "green"),
        ("Upside", 2936, 32, 15, "Agency channel works, ERP/custom deals land", "teal"),
    ]
    for i, (name, revenue, projects, retainers, note, c) in enumerate(scenarios):
        x = 0.72 + i * 2.95
        card(s, x, 1.45, 2.55, 2.0, name, f"{projects} projects | {retainers} retained clients\n{note}", c, f"EGP {revenue:,}K")
        rect(s, x + 0.22, 3.72, 2.1 * revenue / 2936, 0.28, c, c)
    card(s, 0.72, 4.32, 2.55, 0.72, "Monthly GTM burn", "EGP 55K-70K: sales hire, ad spend, agency, tools.", "coral")
    card(s, 3.67, 4.32, 2.55, 0.72, "Break-even logic", "2 core builds or 1 ERP/custom build can cover a lean month.", "gold")
    card(s, 6.62, 4.32, 2.55, 0.72, "Retainer goal", "Build recurring revenue from every successful implementation.", "green")
    footer(s)

    # 12
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Team", "The founding team covers build, growth, and process credibility.", 12)
    people = [
        ("Omar", "AI engineer / full-stack", "Builds demos, automations, integrations, custom systems.", "gold"),
        ("Nabih", "E-commerce / marketing", "Shapes offers, ad angles, customer acquisition, Sprynter proof.", "teal"),
        ("Youssef", "Audit / operations", "Maps workflows, quantifies waste, turns pain into ROI.", "green"),
    ]
    imgs = [ROOT / "public" / "team-omar.jpg", ROOT / "public" / "team-nabih.jpg", ROOT / "public" / "team-youssef.jpg"]
    for i, (name, role, desc, c) in enumerate(people):
        x = 0.72 + i * 3.0
        if imgs[i].exists():
            s.shapes.add_picture(str(imgs[i]), Inches(x + 0.12), Inches(1.55), width=Inches(0.68), height=Inches(0.68))
        card(s, x, 2.32, 2.55, 1.42, role, desc, c, name)
    card(s, 1.0, 4.25, 3.6, 0.68, "Sales hire", "Full-time caller responsible for speed-to-lead, qualification, and booked meetings.", "amber")
    card(s, 5.05, 4.25, 3.6, 0.68, "SM agency partner", "Runs creative, media buying, retargeting, and campaign reporting.", "coral")
    footer(s)

    # 13
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Risks", "Execution risks are real, but controllable with a disciplined delivery model.", 13)
    risks = [
        ("Over-automation", "Keep humans in the loop; automate 60-75% where quality is safe.", "coral"),
        ("Weak client data", "Start with simple flows; clean only the data needed for the first win.", "amber"),
        ("Client resistance", "Demo before payment; train staff and speak in business outcomes.", "teal"),
        ("Scope creep", "Fixed scope, change requests, clear handoff and retainer rules.", "gold"),
        ("Ad quality", "Weekly creative review, lead source tracking, fast budget shifts.", "green"),
        ("Sales execution", "Call SLAs, scripts, CRM hygiene, and founder review of every proposal.", "coral"),
    ]
    for i, (h, b, c) in enumerate(risks):
        x = 0.65 + (i % 3) * 3.0
        y = 1.48 + (i // 3) * 1.48
        card(s, x, y, 2.65, 1.12, h, b, c)
    add_text(s, "Control principle: build the smallest useful system, prove value, then expand.", 0.78, 4.65, 8.45, 0.26, 12, "gold_light", True, align=PP_ALIGN.CENTER, margin=0)
    footer(s)

    # 14
    s = prs.slides.add_slide(blank)
    set_bg(s)
    title(s, "Investor ask", "Funding accelerates the first repeatable sales machine.", 14)
    add_text(s, "Suggested seed ask", 0.72, 1.5, 2.4, 0.22, 11, "muted", True, margin=0)
    add_text(s, "EGP 750K", 0.72, 1.82, 2.9, 0.56, 32, "gold_light", True, font="Aptos Display", margin=0)
    add_text(s, "Six-month runway for ads, sales, demos, and delivery.", 0.75, 2.48, 2.65, 0.5, 10.5, "muted", margin=0)
    uses = [
        ("40%", "Ads + agency", "Demand generation and testing"),
        ("25%", "Sales hire", "Calling, qualification, meetings"),
        ("20%", "Demo templates", "Reusable WhatsApp, CRM, ERP-light systems"),
        ("10%", "Ops/tooling", "CRM, hosting, WhatsApp, workspace"),
        ("5%", "Legal/admin", "Basic contracts and finance hygiene"),
    ]
    for i, (pct, h, b) in enumerate(uses):
        y = 1.25 + i * 0.66
        rect(s, 4.1, y, 0.62, 0.38, ["gold", "teal", "green", "amber", "coral"][i], ["gold", "teal", "green", "amber", "coral"][i])
        add_text(s, pct, 4.12, y + 0.1, 0.58, 0.14, 8, "bg", True, align=PP_ALIGN.CENTER, margin=0)
        add_text(s, h, 4.9, y + 0.02, 1.65, 0.18, 10, "text", True, margin=0)
        add_text(s, b, 6.55, y + 0.03, 2.4, 0.18, 8.5, "muted", margin=0)
    card(s, 0.72, 4.28, 2.58, 0.72, "Milestone gate", "Continue after meetings, paid clients, and retainers are visible.", "green")
    card(s, 3.72, 4.28, 2.58, 0.72, "Investor updates", "Monthly: leads, calls, audits, proposals, revenue, burn.", "teal")
    card(s, 6.72, 4.28, 2.58, 0.72, "Next decision", "Approve launch budget and start agency + sales hiring.", "gold")
    footer(s, "Investor ask is a suggested planning figure and can be adjusted to the actual funding conversation.")

    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    create_deck()
